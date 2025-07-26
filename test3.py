from copy import deepcopy
from io import BytesIO
import io, json, re
from typing import Dict, Any, List
import os
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor

app = FastAPI()

# --------------------- 占位符替换 --------------------- #
PATTERN = re.compile(r"\{\{(\w+)\}\}")

def replace_text_in_runs(para, mapping: Dict[str, Any]):
    """在保留格式的前提下替换文本（逐段处理）"""
    full_text = "".join(run.text for run in para.runs)
    if not PATTERN.search(full_text):
        return False

    new_text = replace_text(full_text, mapping)
    if new_text == full_text:
        return False

    # 关键改进：保留第一段格式，只替换文本内容
    if para.runs:
        # 保留第一个run的格式
        first_run = para.runs[0]
        font = first_run.font

        # 处理字体颜色
        if isinstance(font.color, RGBColor):
            color = font.color.rgb  # 正确访问 .rgb 属性
        else:
            color = None  # 如果是其他颜色类型，跳过

        # 清除所有run
        while para.runs:
            run = para.runs[0]
            if run._r.getparent() == para._p:
                para._p.remove(run._r)
            else:
                break

        # 用保留的格式创建新run
        new_run = para.add_run()
        new_run.text = new_text

        # 复制原始格式
        new_run.font.name = font.name
        new_run.font.size = font.size
        new_run.font.bold = font.bold
        new_run.font.italic = font.italic
        new_run.font.underline = font.underline
        if color:
            new_run.font.color.rgb = color  # 保留原始颜色

    else:
        para.text = new_text
    return True


def replace_text(text: str, mapping: Dict[str, Any]) -> str:
    def repl(m):
        key = m.group(1)
        m2 = re.match(r"^([A-Za-z_]+)(\d+)$", key)
        if m2:
            base, idx = m2.group(1), int(m2.group(2)) - 1
            if base in mapping and isinstance(mapping[base], list) and 0 <= idx < len(mapping[base]):
                return str(mapping[base][idx])
        return str(mapping.get(key, m.group(0)))

    return PATTERN.sub(repl, text)

def _process_group_shape(group, mapping: Dict[str, Any]):
    """递归处理组合形状，保留格式"""
    try:
        for shape in group.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    replace_text_in_runs(para, mapping)

            # 处理组合形状（类型值6）
            if shape.shape_type == 6:
                _process_group_shape(shape, mapping)
    except Exception as e:
        print(f"处理组合形状时出错: {e}")

def process_slide_text(slide, mapping: Dict[str, Any]):
    """处理幻灯片文本，保留原始格式"""
    print(f"\n===== 开始处理幻灯片文本 =====")

    try:
        shape_count = len(slide.shapes)
        print(f"当前形状数量: {shape_count}")
    except Exception as e:
        print(f"获取形状数量时出错: {e}")
        shape_count = 0

    # 刷新形状集合
    if shape_count == 0:
        print("警告：形状集合为空，尝试刷新...")
        try:
            slide.shapes._spTree = slide.element.xpath("./p:cSld/p:spTree")[0]
            shape_count = len(slide.shapes)
            print(f"刷新后形状数量: {shape_count}")
        except Exception as e:
            print(f"刷新形状集合失败: {e}")

    # 处理所有形状
    for i, shape in enumerate(slide.shapes):
        try:
            print(f"\n处理形状 {i + 1}/{len(slide.shapes)}")
            print(f"形状类型值: {shape.shape_type}, 是否有文本框: {shape.has_text_frame}")

            # 处理文本框（保留格式）
            if shape.has_text_frame:
                # 保留文本框的自动调整设置
                tf = shape.text_frame
                autofit = tf.word_wrap  # 保存自动换行设置
                for para in tf.paragraphs:
                    replace_text_in_runs(para, mapping)
                # 恢复自动调整设置
                tf.word_wrap = autofit
                tf.auto_size = tf.auto_size  # 确保自动调整字体大小功能有效

            # 处理组合形状
            if shape.shape_type == 6:
                _process_group_shape(shape, mapping)

        except Exception as e:
            print(f"处理形状 {i + 1} 时出错: {e}")

    print("===== 幻灯片文本处理结束 =====")

# --------------------- 复制工具函数 ------------------- #
def _copy_non_placeholder_shapes(src_cSld, dst_spTree):
    for shp in src_cSld.xpath("./p:spTree/*"):
        if shp.xpath(".//p:ph"):          # 占位符跳过
            continue
        dst_spTree.append(deepcopy(shp))


def _copy_bg(src_cSld, dst_cSld):
    # 尝试获取源幻灯片的背景
    bg = src_cSld.xpath("./p:bg")

    # 打印调试信息：源幻灯片是否包含背景
    print("源幻灯片背景存在:", bool(bg))

    if not bg:
        print("未找到背景，跳过复制操作")
        return

    # 打印调试信息：背景元素的内容
    print("找到的背景内容:", bg)

    # 删除目标幻灯片中的现有背景
    existing_bg = dst_cSld.xpath("./p:bg")
    if existing_bg:
        print("删除目标幻灯片中现有的背景")
        for old in existing_bg:
            dst_cSld.remove(old)
    else:
        print("目标幻灯片没有现有背景，跳过删除操作")

    # 复制背景到目标幻灯片
    dst_cSld.append(deepcopy(bg[0]))
    print("背景已复制到目标幻灯片")


def _copy_all_image_relationships(src_part, dst_slide):
    """
    复制源部件中的所有图片关系到目标幻灯片
    """
    print(f"\n===== 开始复制图片关系 =====")
    
    src_rels = src_part._rels
    dst_part = dst_slide.part
    
    # 获取所有图片关系
    image_rels = []
    for rel_id, rel in src_rels.items():
        if rel.reltype == RT.IMAGE:
            image_rels.append((rel_id, rel))
    
    print(f"源部件中发现 {len(image_rels)} 个图片关系")
    
    # 创建关系ID映射
    rid_mapping = {}
    
    for old_rid, rel in image_rels:
        try:
            # 获取图片数据
            img_blob = rel.target_part.blob
            if not img_blob:
                print(f"图片 {old_rid} 没有数据，跳过")
                continue
            
            # 获取图片内容类型
            content_type = rel.target_part.content_type
            
            print(f"复制图片 {old_rid}，内容类型: {content_type}")
            
            # 创建新的图片部件
            image_part = dst_part.package.get_or_add_image_part(img_blob)
            
            # 在目标幻灯片中添加关系
            new_rid = dst_part._rels.get_or_add_ext_rel(RT.IMAGE, image_part)
            
            rid_mapping[old_rid] = new_rid
            print(f"图片关系映射: {old_rid} -> {new_rid}")
            
        except Exception as e:
            print(f"复制图片关系 {old_rid} 时出错: {e}")
            continue
    
    print(f"成功复制 {len(rid_mapping)} 个图片关系")
    return rid_mapping


def _update_image_references(element, rid_mapping):
    """
    更新元素中的所有图片引用
    """
    print(f"\n===== 开始更新图片引用 =====")
    
    # 查找所有图片引用
    blips = element.xpath(".//a:blip[@r:embed]")
    chart_refs = element.xpath(".//c:externalData[@r:id]")
    ole_refs = element.xpath(".//p:oleObj[@r:id]")
    
    total_refs = len(blips) + len(chart_refs) + len(ole_refs)
    print(f"发现 {total_refs} 个图片引用需要更新")
    
    updated_count = 0
    
    # 更新 a:blip 引用
    for blip in blips:
        old_rid = blip.get(qn("r:embed"))
        if old_rid in rid_mapping:
            new_rid = rid_mapping[old_rid]
            blip.set(qn("r:embed"), new_rid)
            print(f"更新 blip 引用: {old_rid} -> {new_rid}")
            updated_count += 1
    
    # 更新图表引用
    for chart_ref in chart_refs:
        old_rid = chart_ref.get(qn("r:id"))
        if old_rid in rid_mapping:
            new_rid = rid_mapping[old_rid]
            chart_ref.set(qn("r:id"), new_rid)
            print(f"更新图表引用: {old_rid} -> {new_rid}")
            updated_count += 1
    
    # 更新 OLE 对象引用
    for ole_ref in ole_refs:
        old_rid = ole_ref.get(qn("r:id"))
        if old_rid in rid_mapping:
            new_rid = rid_mapping[old_rid]
            ole_ref.set(qn("r:id"), new_rid)
            print(f"更新 OLE 引用: {old_rid} -> {new_rid}")
            updated_count += 1
    
    print(f"成功更新 {updated_count} 个图片引用")


def _fix_blip_rids(src_part, dst_slide):
    """
    修复图片引用：复制所有图片关系并更新引用
    """
    print(f"\n===== 开始修复图片引用 =====")
    print(f"源部件类型: {type(src_part).__name__}")
    
    try:
        # 复制所有图片关系
        rid_mapping = _copy_all_image_relationships(src_part, dst_slide)
        
        if not rid_mapping:
            print("没有图片关系需要复制")
            return
        
        # 更新目标幻灯片中的所有图片引用
        _update_image_references(dst_slide.element, rid_mapping)
        
        print("图片引用修复完成")
        
    except Exception as e:
        print(f"修复图片引用时出错: {e}")
        import traceback
        traceback.print_exc()

# --------------------- 幻灯片克隆 ------------------- #
def clone_slide(dst_prs: Presentation, src_slide):
    """三层克隆：幻灯片自身 + 版式 + 母版"""
    print(f"\n===== 开始克隆幻灯片 =====")
    
    blank_layout = dst_prs.slide_layouts[6]        # 空白
    dst_slide = dst_prs.slides.add_slide(blank_layout)

    dst_cSld   = dst_slide.element.xpath("./p:cSld")[0]
    dst_spTree = dst_cSld.xpath("./p:spTree")[0]
    dst_cSld.remove(dst_spTree)                    # 去掉空白 spTree

    # ----- 幻灯片层 ----- #
    print("复制幻灯片层内容...")
    src_cSld = src_slide.element.xpath("./p:cSld")[0]
    dst_cSld.append(deepcopy(src_cSld.xpath("./p:spTree")[0]))
    _copy_bg(src_cSld, dst_cSld)

    # ----- 版式层 ----- #
    print("复制版式层内容...")
    layout      = src_slide.slide_layout
    layout_cSld = layout.element.xpath("./p:cSld")[0]
    _copy_non_placeholder_shapes(layout_cSld, dst_cSld.xpath("./p:spTree")[0])
    _copy_bg(layout_cSld, dst_cSld)

    # ----- 母版层 ----- #
    print("复制母版层内容...")
    master      = layout.slide_master
    master_cSld = master.element.xpath("./p:cSld")[0]
    _copy_non_placeholder_shapes(master_cSld, dst_cSld.xpath("./p:spTree")[0])
    _copy_bg(master_cSld, dst_cSld)

    # ----- 修正图片关系（按优先级顺序） ----- #
    print("修复图片关系...")
    
    # 1. 首先处理幻灯片本身的图片
    _fix_blip_rids(src_slide.part, dst_slide)
    
    # 2. 然后处理版式层的图片
    _fix_blip_rids(layout.part, dst_slide)
    
    # 3. 最后处理母版层的图片
    _fix_blip_rids(master.part, dst_slide)

    # 返回新幻灯片索引
    slide_index = len(dst_prs.slides) - 1
    print(f"新幻灯片在目标演示文稿中的索引: {slide_index}")
    print("===== 幻灯片克隆结束 =====")
    return slide_index

# --------------------- FastAPI 入口 ------------------ #
@app.post("/merge_and_fill_slides")
async def merge_and_fill_slides(
    templates: List[UploadFile] = File(..., description="多个 PPT 模板"),
    json_data: str              = Form(..., description="合并与替换规则 JSON")
):
    try:
        # 1. 读取模板
        pres_list: List[Presentation] = []
        for tpl in templates:
            pres_list.append(Presentation(BytesIO(await tpl.read())))

        # 2. 解析 JSON
        cfg = json.loads(json_data)
        if "slides" not in cfg:
            raise HTTPException(400, "JSON 必须包含 slides 字段")

        # 3. 创建目标 PPT
        dest = Presentation()
        first = True

        # 4. 逐条配置合并
        for item in cfg["slides"]:
            t_idx = item["template_index"]
            s_idx = item["slide_idx"]
            repl  = item.get("replacements", {})

            if t_idx >= len(pres_list):
                raise HTTPException(400, f"模板索引 {t_idx} 超界")
            src_slide = pres_list[t_idx].slides[s_idx]

            # 同步尺寸
            if first:
                dest.slide_width  = pres_list[t_idx].slide_width
                dest.slide_height = pres_list[t_idx].slide_height
                first = False

            slide_index = clone_slide(dest, src_slide)
            # 获取新幻灯片并替换文本
            new_slide = dest.slides[slide_index]
            if repl:
                print(f"开始替换占位符，替换数据: {repl}")
                process_slide_text(new_slide, repl)
            else:
                print("无替换数据，跳过替换步骤")

        # 5. 输出文件
        buf = io.BytesIO()
        dest.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=merged.pptx"}
        )

    except json.JSONDecodeError as e:
        raise HTTPException(400, f"JSON 解析失败：{e}")
    except Exception as e:
        raise HTTPException(500, f"处理失败：{e}")